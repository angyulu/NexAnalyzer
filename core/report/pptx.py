"""
Assembles the Sample Report's three-slide PPTX entirely in memory.

Layout (16:9 slides):
- Slide 1 (overview): title bar; 3x3 OM image grid on the left half;
  Raman fit-summary table (optionally with a trailing amplitude-ratio row)
  stacked above the PL fit-summary table on the right half.
- Slide 2: the nine Raman points as three stacked columns (1/4/7, 2/5/8,
  3/6/9), each column one image whose three panels share an X-axis, under
  one legend and one pair of axis titles serving the whole grid.
- Slide 3: same, for PL.
"""

import io
from typing import Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .models import PeakStat

SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN = 13.333, 7.5

_TITLE_LEFT, _TITLE_TOP, _TITLE_W, _TITLE_H = 0.30, 0.15, 12.73, 0.65
_RULE_TOP, _RULE_H = 0.83, 0.02
_CONTENT_LEFT, _CONTENT_TOP, _CONTENT_W = 0.30, 0.95, 12.73

# ---- Slide 1: OM grid (left half) ----
_GRID_W, _GRID_H = 6.20, 4.55
_GRID_GUTTER = 0.10
_GRID_CAPTION_TOP, _GRID_CAPTION_H = 5.52, 0.28

# ---- Slide 1: stats tables (right half, stacked) ----
_TABLE_LEFT = 6.70
_TABLE_W = 6.33
_TABLE_H = 2.45
_TABLE_CAPTION_GAP = 0.27
_RAMAN_TABLE_TOP = 1.22
# The Raman table gets the extra height because it may carry an additional
# amplitude-ratio row; PL never does.
_RAMAN_TABLE_H = 2.73
_PL_TABLE_TOP = 4.27

# ---- Slides 2 & 3: 3x3 fitted-spectrum grid, with one shared legend ----
# The legend and axis titles live on the slide rather than inside each cell
# image. Nine copies of the same key and the same two axis titles were most of
# what made the individual spectra small; drawn once, the cells keep that space.
_FIT_LEGEND_TOP, _FIT_LEGEND_H = 0.90, 0.24
_FIT_YLABEL_W = 0.26  # left gutter holding the rotated Y-axis title
_FIT_XLABEL_TOP, _FIT_XLABEL_H = 7.12, 0.28

_FIT_GRID_LEFT = _CONTENT_LEFT + _FIT_YLABEL_W
_FIT_GRID_TOP = 1.18
_FIT_GRID_W = SLIDE_WIDTH_IN - _FIT_GRID_LEFT - _CONTENT_LEFT
_FIT_GRID_H = 5.92
# Columns sit close together too: the gutter only has to separate them, and
# every hundredth of an inch here goes to the spectra.
_FIT_GRID_GUTTER = 0.06

_LEGEND_SWATCH_W, _LEGEND_SWATCH_H, _LEGEND_GAP = 0.20, 0.05, 0.06
_LEGEND_FONT, _AXIS_LABEL_FONT = Pt(12), Pt(12)

# Three columns across the grid's width, each running its full height: a
# column is one image holding three stacked panels that share an X-axis.
FIT_GRID_COLUMNS = 3
_FIT_COLUMN_W = (_FIT_GRID_W - (FIT_GRID_COLUMNS - 1) * _FIT_GRID_GUTTER) / FIT_GRID_COLUMNS

# Public: callers rendering a column figure to PNG (e.g. via
# export_figure_png) should match this aspect ratio, otherwise
# aspect-preserving _fit_picture() letterboxes the image inside its column
# and leaves part of it empty.
FIT_COLUMN_ASPECT_RATIO = _FIT_COLUMN_W / _FIT_GRID_H

_PLACEHOLDER_LINE = RGBColor(0x00, 0x00, 0x00)


def build_sample_report_pptx(
    sample_name: str,
    material_name: str,
    report_date: str,
    magnification_label: Optional[str],
    om_image_bytes: Dict[int, bytes],
    raman_stats: Optional[List[PeakStat]],
    raman_fit_columns: Dict[int, bytes],
    pl_stats: Optional[List[PeakStat]],
    pl_fit_columns: Dict[int, bytes],
    raman_amplitude_ratio: Optional[Tuple[float, float, int]] = None,
    raman_amplitude_ratio_label: str = "",
    raman_fit_legend: Optional[List[Tuple[str, str]]] = None,
    pl_fit_legend: Optional[List[Tuple[str, str]]] = None,
    raman_x_label: str = "Raman Shift (cm⁻¹)",
    pl_x_label: str = "Wavelength (nm)",
    fit_y_label: str = "Normalized intensity",
) -> bytes:
    """Build the three-slide sample report and return .pptx bytes.

    `raman_amplitude_ratio`, if given, is (median, MAD, n) of a per-point
    peak-height ratio (e.g. LA/E2g+A1g for WSe2 Raman — see
    peak_metrics.compute_peak_height_ratio) appended as a final, bolded row
    of the Raman fit-summary table; omitted entirely when None.

    `raman_fit_columns`/`pl_fit_columns` map column index (0, 1, 2) to one PNG
    holding that column's three stacked points — 0 is points 1/4/7, 1 is 2/5/8,
    2 is 3/6/9 — rendered at FIT_COLUMN_ASPECT_RATIO. A missing column becomes
    a placeholder.

    `raman_fit_legend`/`pl_fit_legend` are (label, "#RRGGBB") pairs describing
    the traces in that technique's grid, drawn once above it. The column images
    are expected to carry no legend of their own; passing None just omits it.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    overview_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(overview_slide, sample_name, material_name, report_date)
    _add_om_grid(overview_slide, om_image_bytes, magnification_label)
    _add_stats_table(
        overview_slide, "Raman", raman_stats, _TABLE_LEFT, _RAMAN_TABLE_TOP, _TABLE_W, _RAMAN_TABLE_H,
        ratio=raman_amplitude_ratio, ratio_label=raman_amplitude_ratio_label,
    )
    _add_stats_table(overview_slide, "PL", pl_stats, _TABLE_LEFT, _PL_TABLE_TOP, _TABLE_W, _TABLE_H)

    raman_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(raman_slide, sample_name, material_name, report_date, subtitle="Raman — fitted spectra (9 points)")
    _add_fit_legend(raman_slide, raman_fit_legend)
    _add_fit_axis_titles(raman_slide, raman_x_label, fit_y_label)
    _add_fit_columns(raman_slide, raman_fit_columns)

    pl_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(pl_slide, sample_name, material_name, report_date, subtitle="PL — fitted spectra (9 points)")
    _add_fit_legend(pl_slide, pl_fit_legend)
    _add_fit_axis_titles(pl_slide, pl_x_label, fit_y_label)
    _add_fit_columns(pl_slide, pl_fit_columns)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_bar(
    slide, sample_name: str, material_name: str, report_date: str, subtitle: Optional[str] = None
) -> None:
    box = slide.shapes.add_textbox(Inches(_TITLE_LEFT), Inches(_TITLE_TOP), Inches(_TITLE_W), Inches(_TITLE_H))
    tf = box.text_frame
    text = f"{sample_name}   |   Material: {material_name}   |   {report_date}"
    if subtitle:
        text += f"   |   {subtitle}"
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(_TITLE_LEFT), Inches(_RULE_TOP), Inches(_TITLE_W), Inches(_RULE_H)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x40)
    rule.line.fill.background()


def _add_om_grid(slide, om_image_bytes: Dict[int, bytes], magnification_label: Optional[str]) -> None:
    cell_w = (_GRID_W - 2 * _GRID_GUTTER) / 3
    cell_h = (_GRID_H - 2 * _GRID_GUTTER) / 3

    for point in range(1, 10):
        row, col = divmod(point - 1, 3)
        cell_left = _CONTENT_LEFT + col * (cell_w + _GRID_GUTTER)
        cell_top = _CONTENT_TOP + row * (cell_h + _GRID_GUTTER)

        image_bytes = om_image_bytes.get(point)
        if image_bytes is not None:
            _fit_picture(slide, image_bytes, cell_left, cell_top, cell_w, cell_h)
        else:
            _add_placeholder(slide, cell_left, cell_top, cell_w, cell_h)

    caption_text = "OM"
    if magnification_label:
        caption_text += f" ({magnification_label})"
    caption = slide.shapes.add_textbox(
        Inches(_CONTENT_LEFT), Inches(_GRID_CAPTION_TOP), Inches(_GRID_W), Inches(_GRID_CAPTION_H)
    )
    caption.text_frame.text = caption_text
    caption.text_frame.paragraphs[0].font.size = Pt(11)
    caption.text_frame.paragraphs[0].font.italic = True


def _hex_to_rgb(value: str) -> RGBColor:
    """A peak template's "#RRGGBB" as a pptx color, falling back to black for
    anything unparseable (FittedPeak.color, unlike PeakDefinition.color, is
    never validated)."""
    try:
        raw = value.lstrip("#")
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except (AttributeError, ValueError, IndexError):
        return RGBColor(0x00, 0x00, 0x00)


def _add_fit_legend(slide, entries: Optional[List[Tuple[str, str]]]) -> None:
    """One legend for all nine cells: a colored swatch and label per trace,
    spread evenly across the grid's width above it.

    Drawn as shapes rather than rendered into the cell images so it appears
    once instead of nine times, and stays crisp at any zoom.
    """
    if not entries:
        return

    slot_w = _FIT_GRID_W / len(entries)
    swatch_top = _FIT_LEGEND_TOP + (_FIT_LEGEND_H - _LEGEND_SWATCH_H) / 2

    for i, (label, color) in enumerate(entries):
        slot_left = _FIT_GRID_LEFT + i * slot_w

        swatch = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(slot_left), Inches(swatch_top),
            Inches(_LEGEND_SWATCH_W), Inches(_LEGEND_SWATCH_H),
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = _hex_to_rgb(color)
        swatch.line.fill.background()

        text_left = slot_left + _LEGEND_SWATCH_W + _LEGEND_GAP
        box = slide.shapes.add_textbox(
            Inches(text_left), Inches(_FIT_LEGEND_TOP),
            Inches(max(0.1, slot_w - _LEGEND_SWATCH_W - _LEGEND_GAP)), Inches(_FIT_LEGEND_H),
        )
        tf = box.text_frame
        tf.text = label
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].font.size = _LEGEND_FONT


def _add_fit_axis_titles(slide, x_label: str, y_label: str) -> None:
    """The axis titles the column images no longer carry, written once each:
    the Y title rotated in the left gutter, the X title centered underneath."""
    y_box = slide.shapes.add_textbox(
        Inches(_CONTENT_LEFT), Inches(_FIT_GRID_TOP), Inches(_FIT_YLABEL_W), Inches(_FIT_GRID_H)
    )
    y_tf = y_box.text_frame
    y_tf.text = y_label
    y_tf.word_wrap = False
    y_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    y_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_tf.paragraphs[0].font.size = _AXIS_LABEL_FONT
    # Rotate the text inside the box rather than the box itself: a rotated
    # shape is positioned about its own center, which puts its top-left well
    # off the slide for a label this tall.
    y_tf._txBody.find(qn("a:bodyPr")).set("vert", "vert270")

    x_box = slide.shapes.add_textbox(
        Inches(_FIT_GRID_LEFT), Inches(_FIT_XLABEL_TOP), Inches(_FIT_GRID_W), Inches(_FIT_XLABEL_H)
    )
    x_tf = x_box.text_frame
    x_tf.text = x_label
    x_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    x_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_tf.paragraphs[0].font.size = _AXIS_LABEL_FONT


def _add_fit_columns(slide, column_images: Dict[int, bytes]) -> None:
    """The grid, as three full-height columns side by side.

    A column arrives as one image because its three points share an X-axis,
    which only holds if they are drawn in a single figure. A column with no
    image (an omitted technique, or every point in it failing to fit) renders
    as a placeholder, keeping the other columns in their usual positions. The
    legend and axis titles are added separately, once per slide."""
    for col in range(FIT_GRID_COLUMNS):
        col_left = _FIT_GRID_LEFT + col * (_FIT_COLUMN_W + _FIT_GRID_GUTTER)

        image_bytes = column_images.get(col)
        if image_bytes is not None:
            _fit_picture(slide, image_bytes, col_left, _FIT_GRID_TOP, _FIT_COLUMN_W, _FIT_GRID_H)
        else:
            _add_placeholder(slide, col_left, _FIT_GRID_TOP, _FIT_COLUMN_W, _FIT_GRID_H)


def _add_stats_table(
    slide, technique_label: str, stats: Optional[List[PeakStat]], left: float, top: float, w: float, h: float,
    ratio: Optional[Tuple[float, float, int]] = None, ratio_label: str = "",
) -> None:
    """One technique's fit-summary table.

    `ratio`, if given, is appended as a final bolded row: its label in the
    Peak column and `median ± MAD` in the Amplitude column (it is a ratio of
    heights), with center/FWHM dashed out since they don't apply. The label is
    expected to mark it as a median, since the peak rows above are means.
    """
    caption = slide.shapes.add_textbox(
        Inches(left), Inches(top - _TABLE_CAPTION_GAP), Inches(w), Inches(_TABLE_CAPTION_GAP)
    )
    caption.text_frame.text = f"{technique_label} fit summary (mean ± std)"
    caption.text_frame.paragraphs[0].font.size = Pt(12)
    caption.text_frame.paragraphs[0].font.bold = True

    if not stats:
        _add_placeholder(slide, left, top, w, h)
        return

    n_data_rows = len(stats) + (1 if ratio is not None else 0)
    n_rows = n_data_rows + 1
    font_size = Pt(13) if n_data_rows <= 6 else (Pt(11) if n_data_rows <= 10 else Pt(9))

    table_shape = slide.shapes.add_table(n_rows, 5, Inches(left), Inches(top), Inches(w), Inches(h))
    table = table_shape.table

    headers = ["Peak", "Center", "Amplitude", "FWHM", "n"]
    col_fracs = [0.22, 0.26, 0.26, 0.18, 0.08]
    for c, frac in enumerate(col_fracs):
        table.columns[c].width = Inches(w * frac)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = font_size
        cell.text_frame.paragraphs[0].font.bold = True

    for r, stat in enumerate(stats, start=1):
        values = [
            stat.label,
            f"{stat.center_mean:.1f} ± {stat.center_std:.1f}",
            f"{stat.height_mean:.1f} ± {stat.height_std:.1f}",
            f"{stat.fwhm_mean:.1f} ± {stat.fwhm_std:.1f}",
            str(stat.n),
        ]
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = font_size

    if ratio is not None:
        median, mad, n = ratio
        # Three decimals: these ratios run around 0.1, where two would round
        # away most of the point-to-point variation the row exists to show.
        values = [ratio_label, "—", f"{median:.3f} ± {mad:.3f}", "—", str(n)]
        for c, value in enumerate(values):
            cell = table.cell(n_rows - 1, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = font_size
            cell.text_frame.paragraphs[0].font.bold = True


_EMBED_DPI = 200  # target resolution for embedded pictures at their displayed size


def _fit_picture(slide, image_bytes: bytes, left: float, top: float, max_w: float, max_h: float) -> None:
    """Add `image_bytes` to `slide`, aspect-preserving-scaled to fit within
    (max_w, max_h) inches, centered in that box.

    Downsamples source images whose native resolution exceeds _EMBED_DPI at
    their displayed size (e.g. multi-megapixel microscope photos going into
    a postage-stamp-sized grid cell) so the .pptx doesn't balloon from
    embedding full-resolution originals; images already at or below that
    resolution are embedded unchanged.
    """
    with Image.open(io.BytesIO(image_bytes)) as im:
        native_w, native_h = im.size
        native_ratio = native_w / native_h
        box_ratio = max_w / max_h

        if native_ratio > box_ratio:
            width, height = max_w, max_w / native_ratio
        else:
            width, height = max_h * native_ratio, max_h

        target_px = (max(1, round(width * _EMBED_DPI)), max(1, round(height * _EMBED_DPI)))
        if native_w > target_px[0] or native_h > target_px[1]:
            im.thumbnail(target_px, Image.LANCZOS)
            out_buf = io.BytesIO()
            im.save(out_buf, format="PNG")
            image_bytes = out_buf.getvalue()

    left_offset = left + (max_w - width) / 2
    top_offset = top + (max_h - height) / 2

    slide.shapes.add_picture(
        io.BytesIO(image_bytes), Inches(left_offset), Inches(top_offset),
        width=Inches(width), height=Inches(height)
    )


def _add_placeholder(slide, left: float, top: float, w: float, h: float):
    """An empty black-outlined, unfilled box marking missing content (no
    image, no fit, no stats) — deliberately blank rather than explaining
    why, so the report stays visually clean."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = _PLACEHOLDER_LINE
    shape.line.width = Pt(1)
    return shape
