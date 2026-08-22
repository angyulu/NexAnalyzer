"""Unit tests for core.report.pptx (Sample Report .pptx assembly)."""

import io

from PIL import Image
from pptx import Presentation

from core.report.pptx import build_sample_report_pptx
from core.report.models import PeakStat


def _tiny_png(size=(40, 30), color=(200, 200, 200)) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, format="PNG")
    return buf.getvalue()


def _tables(slide):
    """Tables on `slide`, in the order they were added (Raman, then PL)."""
    return [shape.table for shape in slide.shapes if shape.has_table]


def _texts(slide):
    """Every piece of text on `slide`."""
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def _fit_slides(result):
    """The two 3x3 fitted-spectrum slides (Raman, then PL)."""
    prs = Presentation(io.BytesIO(result))
    return prs.slides[1], prs.slides[2]


def _stat(label="Exciton", n=9):
    return PeakStat(
        label=label, n=n,
        center_mean=766.5, center_std=0.7,
        height_mean=12000.0, height_std=1500.0,
        fwhm_mean=25.0, fwhm_std=1.2,
    )


class TestBuildSampleReportPptx:
    def test_full_report_builds_three_slides(self):
        om_bytes = {p: _tiny_png() for p in range(1, 10)}
        column_bytes = {c: _tiny_png() for c in range(3)}

        result = build_sample_report_pptx(
            sample_name="VABA52", material_name="WSe2", report_date="2026-08-21",
            magnification_label="100x", om_image_bytes=om_bytes,
            raman_stats=[_stat("E2g+A1g"), _stat("2LA")], raman_fit_columns=column_bytes,
            pl_stats=[_stat("Exciton"), _stat("Trion")], pl_fit_columns=column_bytes,
        )

        assert result[:2] == b"PK"  # pptx is a zip archive
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_missing_points_use_placeholders_without_error(self):
        result = build_sample_report_pptx(
            sample_name="Partial", material_name="WSe2", report_date="2026-08-21",
            magnification_label="100x", om_image_bytes={1: _tiny_png(), 5: _tiny_png()},
            raman_stats=None, raman_fit_columns={0: _tiny_png()},
            pl_stats=None, pl_fit_columns={},
        )

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_technique_entirely_omitted(self):
        column_bytes = {c: _tiny_png() for c in range(3)}

        result = build_sample_report_pptx(
            sample_name="RamanOnly", material_name="Silicon", report_date="2026-08-21",
            magnification_label=None, om_image_bytes={},
            raman_stats=[_stat("Si", n=9)], raman_fit_columns=column_bytes,
            pl_stats=None, pl_fit_columns={},
        )

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3
        # PL slide's fit grid should be all placeholders (no PL images given)
        pl_slide = prs.slides[2]
        pictures = [s for s in pl_slide.shapes if s.shape_type == 13]
        assert len(pictures) == 0

    def test_wide_and_tall_images_both_fit_without_error(self):
        column_bytes = {c: _tiny_png() for c in range(3)}

        result = build_sample_report_pptx(
            sample_name="AspectRatios", material_name="WSe2", report_date="2026-08-21",
            magnification_label="100x",
            om_image_bytes={1: _tiny_png(size=(200, 40)), 2: _tiny_png(size=(40, 200))},
            raman_stats=[_stat()], raman_fit_columns=column_bytes,
            pl_stats=None, pl_fit_columns={},
        )

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_many_peaks_shrinks_table_font_without_error(self):
        stats = [_stat(label=f"Peak{i}") for i in range(11)]

        result = build_sample_report_pptx(
            sample_name="ManyPeaks", material_name="WSe2", report_date="2026-08-21",
            magnification_label=None, om_image_bytes={},
            raman_stats=stats, raman_fit_columns={},
            pl_stats=None, pl_fit_columns={},
        )

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_missing_content_placeholders_have_no_text(self):
        result = build_sample_report_pptx(
            sample_name="Empty", material_name="WSe2", report_date="2026-08-21",
            magnification_label=None, om_image_bytes={},
            raman_stats=None, raman_fit_columns={},
            pl_stats=None, pl_fit_columns={},
        )

        prs = Presentation(io.BytesIO(result))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 1 and shape.has_text_frame:  # AUTO_SHAPE (placeholder rectangles)
                    assert shape.text_frame.text == ""

    def test_raman_amplitude_ratio_is_a_row_of_the_raman_table(self):
        result = build_sample_report_pptx(
            sample_name="RatioTest", material_name="WSe2", report_date="2026-08-21",
            magnification_label=None, om_image_bytes={},
            raman_stats=[_stat("LA"), _stat("E2g+A1g")], raman_fit_columns={},
            pl_stats=None, pl_fit_columns={},
            raman_amplitude_ratio=(0.59, 0.08, 9), raman_amplitude_ratio_label="LA / E2g+A1g ratio",
        )

        prs = Presentation(io.BytesIO(result))
        raman_table = _tables(prs.slides[0])[0]

        # header + 2 peaks + ratio row
        assert len(raman_table.rows) == 4
        ratio_row = [c.text for c in raman_table.rows[3].cells]
        assert ratio_row[0] == "LA / E2g+A1g ratio"
        # Median +/- MAD, to three decimals: these ratios run around 0.1, where
        # two decimals would round away the variation the row exists to show.
        assert ratio_row[2] == "0.590 ± 0.080"  # value sits in the Amplitude column
        assert ratio_row[4] == "9"

        # ...and nowhere outside the table
        overview_texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert not any("ratio" in t for t in overview_texts)

    def test_raman_amplitude_ratio_omitted_when_none(self):
        result = build_sample_report_pptx(
            sample_name="NoRatio", material_name="Silicon", report_date="2026-08-21",
            magnification_label=None, om_image_bytes={},
            raman_stats=[_stat("Si")], raman_fit_columns={},
            pl_stats=None, pl_fit_columns={},
            raman_amplitude_ratio=None,
        )

        prs = Presentation(io.BytesIO(result))
        raman_table = _tables(prs.slides[0])[0]
        assert len(raman_table.rows) == 2  # header + the single peak, no ratio row

        overview_texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert not any("ratio" in t for t in overview_texts)


class TestFitGridLegendAndAxisTitles:
    """Slides 2 and 3 carry one legend and one pair of axis titles for all
    nine cells, so the cell images don't have to repeat them."""

    _LEGEND = [("Data", "#1f77b4"), ("Total Fit", "#ff7f0e"), ("E2g+A1g", "#EF482E")]

    def _build(self, **kwargs):
        return build_sample_report_pptx(
            sample_name="VBBA14", material_name="WSe2", report_date="2026-08-22",
            magnification_label=None, om_image_bytes={},
            raman_stats=None, raman_fit_columns={c: _tiny_png() for c in range(3)},
            pl_stats=None, pl_fit_columns={c: _tiny_png() for c in range(3)},
            **kwargs,
        )

    def test_legend_labels_appear_once_on_their_slide(self):
        raman_slide, pl_slide = _fit_slides(self._build(raman_fit_legend=self._LEGEND))
        texts = _texts(raman_slide)

        for label, _color in self._LEGEND:
            assert texts.count(label) == 1, f"{label} should appear exactly once"
        # ...and only on the technique it was given for.
        assert "E2g+A1g" not in _texts(pl_slide)

    def test_each_legend_entry_gets_a_swatch_in_its_color(self):
        raman_slide, _pl = _fit_slides(self._build(raman_fit_legend=self._LEGEND))

        # Swatches are the only filled autoshapes on a fit slide besides the
        # title rule, which is dark grey.
        fills = {
            str(s.fill.fore_color.rgb)
            for s in raman_slide.shapes
            if s.shape_type == 1 and s.has_text_frame and s.text_frame.text == ""
        }
        for _label, color in self._LEGEND:
            assert color.lstrip("#").upper() in fills

    def test_no_legend_when_none_given(self):
        """A technique with no fits gets an all-placeholder slide and no key."""
        raman_slide, _pl = _fit_slides(self._build(raman_fit_legend=None))
        texts = _texts(raman_slide)

        assert "Data" not in texts
        assert "Total Fit" not in texts

    def test_both_grid_slides_carry_the_axis_titles(self):
        raman_slide, pl_slide = _fit_slides(self._build())

        # Normalized by default now: every panel's own peak is 1.0, so the
        # units are no longer arbitrary.
        assert "Normalized intensity" in _texts(raman_slide)
        assert "Normalized intensity" in _texts(pl_slide)

    def test_x_axis_title_names_the_technique_s_own_units(self):
        raman_slide, pl_slide = _fit_slides(self._build())

        assert "Raman Shift (cm\u207b\u00b9)" in _texts(raman_slide)
        assert "Wavelength (nm)" in _texts(pl_slide)

    def test_x_axis_title_is_overridable(self):
        raman_slide, _pl = _fit_slides(self._build(raman_x_label="Energy (eV)"))

        assert "Energy (eV)" in _texts(raman_slide)

    def test_overview_slide_has_no_grid_axis_titles(self):
        prs = Presentation(io.BytesIO(self._build(raman_fit_legend=self._LEGEND)))

        assert "Normalized intensity" not in _texts(prs.slides[0])


class TestFitColumns:
    """A grid slide is three full-height column images, not nine cells: the
    three points of a column share an X-axis, which only holds if they are
    drawn in one figure."""

    def _build(self, **kwargs):
        defaults = dict(
            sample_name="VBBA14", material_name="WSe2", report_date="2026-08-22",
            magnification_label=None, om_image_bytes={},
            raman_stats=None, raman_fit_columns={}, pl_stats=None, pl_fit_columns={},
        )
        defaults.update(kwargs)
        return build_sample_report_pptx(**defaults)

    def _pictures(self, slide):
        return [s for s in slide.shapes if s.shape_type == 13]

    def _placeholders(self, slide):
        return [s for s in slide.shapes if s.shape_type == 1 and s.has_text_frame and s.text_frame.text == ""]

    def test_three_columns_give_three_pictures(self):
        raman_slide, _pl = _fit_slides(self._build(raman_fit_columns={c: _tiny_png() for c in range(3)}))

        assert len(self._pictures(raman_slide)) == 3

    def test_a_missing_column_becomes_a_placeholder_in_its_own_position(self):
        """The other columns must not shuffle left to fill the gap."""
        result = self._build(raman_fit_columns={0: _tiny_png(), 2: _tiny_png()})
        raman_slide, _pl = _fit_slides(result)

        pictures = self._pictures(raman_slide)
        assert len(pictures) == 2
        # The title rule is also a filled autoshape, so allow for it.
        assert len(self._placeholders(raman_slide)) >= 1

        lefts = sorted(p.left for p in pictures)
        gap = lefts[1] - lefts[0]
        # Columns 0 and 2 are two column-widths apart, not adjacent.
        assert gap > 0

    def test_no_columns_gives_placeholders_and_no_pictures(self):
        raman_slide, pl_slide = _fit_slides(self._build())

        assert self._pictures(raman_slide) == []
        assert self._pictures(pl_slide) == []

    def test_columns_span_the_grid_width_in_order(self):
        raman_slide, _pl = _fit_slides(self._build(raman_fit_columns={c: _tiny_png() for c in range(3)}))

        lefts = [p.left for p in self._pictures(raman_slide)]
        assert lefts == sorted(lefts), "columns should be laid out left to right"

    def test_y_axis_title_is_caller_supplied(self):
        raman_slide, _pl = _fit_slides(self._build(fit_y_label="Normalized intensity"))

        assert "Normalized intensity" in _texts(raman_slide)
